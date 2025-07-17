#  实际可行

#  获取幻灯片中的所有文字

import collections.abc   #  防止发生异常
from pptx import Presentation

pre1 = Presentation('pptx_Trans\\RT-Thread.pptx')

slide_list = pre1.slides

for s_idx, slide in enumerate(slide_list):
    print(f'\n===== 第{s_idx+1}页幻灯片 =====')
    shape_list = slide.shapes
    print('方框对象为')
    for shape in shape_list:
        print(shape)

    # 遍历所有有 text_frame 的 shape
    for idx, shape in enumerate(shape_list):
        if hasattr(shape, "text_frame") and shape.has_text_frame:
            text_frame = getattr(shape, "text_frame")
            print(f'\n  第{idx+1}个文本框对象为:', text_frame)

            para_list = text_frame.paragraphs
            print('  段落对象为:')
            for p_idx, para in enumerate(para_list):
                print(f'    第{p_idx+1}段: {para}')

                run_list = para.runs
                print('    文字块对象为:')
                for r_idx, run in enumerate(run_list):
                    print(f'      第{r_idx+1}块: {run.text}')

