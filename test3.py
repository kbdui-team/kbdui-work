#  pptx 转 docx

import collections.abc
from docx import Document
from pptx import Presentation

def extract_ppt_to_docx(ppt_path, docx_path):
    """
    将PowerPoint文件内容提取到Word文档
    
    参数:
        ppt_path: PowerPoint文件路径
        docx_path: 输出的Word文档路径
    """
    # 读取PPT文件
    ppt = Presentation(ppt_path)
    # 创建Word文档
    doc = Document()
    
    # 添加标题
    doc.add_heading(f'PPT内容提取: {ppt_path}', level=1)
    
    # 遍历所有幻灯片
    for i, slide in enumerate(ppt.slides, start=1):
        # 添加幻灯片标题
        doc.add_heading(f'幻灯片 {i}', level=2)
        
        # 遍历幻灯片中的所有形状
        for shape in slide.shapes:
            # 检查形状是否有文本框
            if shape.has_text_frame:
                text_frame = getattr(shape, "text_frame")
                
                # 遍历文本框中的所有段落
                for paragraph in text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:  # 忽略空段落
                        doc.add_paragraph(text)
        
        # 幻灯片之间添加分隔线
        doc.add_paragraph("-" * 50)
    
    # 保存Word文档
    doc.save(docx_path)
    print(f"成功提取内容到: {docx_path}")

if __name__ == "__main__":
    # 使用示例
    pptx_file = 'pptx_Trans\\RT-Thread.pptx'
    docx_file = 'pptx_Trans\\RT-Thread.docx'
    extract_ppt_to_docx(pptx_file, docx_file)
